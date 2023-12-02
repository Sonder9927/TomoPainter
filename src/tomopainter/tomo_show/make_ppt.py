from pathlib import Path

from pptx import Presentation

from .add_slide import (
    ppt_add_dcs,
    ppt_add_probs,
    ppt_add_prob_with_dc,
    ppt_add_profile,
    ppt_add_one_fig_per_slide,
    ppt_add_single_type,
)


class PptMaker:
    def __init__(self, pn, fig_root: Path, remake=True) -> None:
        self.ppt_name = pn
        self.figs = fig_root
        self.margins = [[4.567, 2.345], [0.789, 0.567]]
        self.shape = {
            "center": [15, 14.5],
            "sub1": [7, 8.5],
            "sub2": [11.3, 5.5],
            "dc": [7.125, 3.33],
            "prob": [5.7, 9],
        }
        # create a new instance for pptx
        self.prs = self._ppt(remake)

    def add_area(self, area_dir):
        area = self.figs / area_dir
        self.prs = ppt_add_one_fig_per_slide(self.prs, area, [3.13, 2.45],
                                             [20, 11], "area.png")
        self.prs = ppt_add_one_fig_per_slide(self.prs, area, [4.56, 2.1],
                                             [14.5, 15.5], "evt_sites.png")
        self.prs = ppt_add_one_fig_per_slide(self.prs, area, [3.456, 2.345],
                                             [18, 10], "perNum_vel.png")
        # self.prs = ppt_add_one_fig_per_slide(
        #     self.prs, area, [3.6, 2.1], [15, 15], "rays_cover.png"
        # )
        self.prs = ppt_add_one_fig_per_slide(self.prs, area, [1.5, 2.45],
                                             [21, 12], "model.png")

    def add_phase_results(self, phase_dir):
        phase = self.figs / phase_dir
        # ant
        self.prs = ppt_add_single_type(
            self.prs,
            phase / "ant_figs",
            "*VEL*",
            self.margins[1],
            self.shape["sub1"],
            key=lambda p: int(p.stem.split("_")[-1]),
        )
        # tpwt
        tpwt = phase / "tpwt_figs"
        # add phase vel of all periods
        self.prs = ppt_add_single_type(
            self.prs,
            tpwt / "phv",
            "*VEL*",
            self.margins[1],
            self.shape["sub1"],
            key=lambda p: int(p.stem.split("_")[-1]),
        )
        # ave vel and std
        self.prs = ppt_add_single_type(
            self.prs,
            tpwt / "as",
            "*AS*",
            self.margins[1],
            self.shape["sub2"],
            key=lambda p: int(p.stem.split("_")[-1]),
            rcn=[3, 2],
        )
        # add check board of all periods
        self.prs = ppt_add_single_type(
            self.prs,
            tpwt / "checkboard",
            "*CB*",
            self.margins[1],
            self.shape["sub1"],
            key=lambda p: int(p.stem.split("_")[-1]),
        )
        # add diff between ant and tpwt
        self.prs = ppt_add_one_fig_per_slide(
            self.prs,
            phase / "diff_figs",
            self.margins[0],
            self.shape["center"],
            "diff*",
        )

    def add_mc_results(self, mc_dir):
        mc = self.figs / mc_dir
        # probalCrs
        self.prs = ppt_add_probs(
            self.prs, mc / "prob", self.margins[1], self.shape["prob"]
        )
        # misfit
        self.prs = ppt_add_one_fig_per_slide(
            self.prs,
            mc,
            self.margins[0],
            [15, 12.5],
            "misfit.png",
        )
        # bad grid
        self.prs = ppt_add_prob_with_dc(
            self.prs,
            mc / "bad",
            [4.567, 1.234],
            [self.shape["dc"], self.shape["prob"]],
        )
        # vs depth
        self.prs = ppt_add_single_type(
            self.prs,
            mc / "depth",
            "vs*",
            self.margins[1],
            self.shape["sub1"],
            key=lambda p: int(p.stem.split("_")[-1][:-2]),
        )
        # vs profile
        for lt in (mc / "profile").glob("*"):
            self.prs = ppt_add_profile(self.prs, lt, [0.35, 1.83],
                                       self.shape["center"])

    def add_dispersion_curves(self, dcs_dir):
        """
        add diff of all periods
        """
        self.prs = ppt_add_dcs(self.prs, self.figs / dcs_dir, self.margins[1],
                               self.shape["dc"])

    def save(self, target=None):
        """
        save ppt
        """
        pn = target or self.ppt_name
        self.prs.save(pn)

    def _ppt(self, remake):
        prs = Presentation() if remake else Presentation(self.ppt_name)
        # make first slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Results Show"
        return prs
